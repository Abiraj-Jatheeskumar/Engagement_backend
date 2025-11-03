"""
Text extraction from PDF and PPTX files
"""
import fitz  # PyMuPDF
from pptx import Presentation
from typing import List
import io


def extract_text_from_pdf(file_content: bytes) -> List[str]:
    """
    Extract text from PDF file
    
    Args:
        file_content: PDF file as bytes
    
    Returns:
        List of text chunks from each page
    """
    try:
        doc = fitz.open(stream=io.BytesIO(file_content), filetype="pdf")
        text_chunks = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            # Clean and chunk text
            cleaned_text = clean_text(text)
            if cleaned_text:
                text_chunks.append(cleaned_text)
        
        doc.close()
        return text_chunks
    except Exception as e:
        raise Exception(f"Error extracting text from PDF: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> List[str]:
    """
    Extract text from PPTX file
    
    Args:
        file_content: PPTX file as bytes
    
    Returns:
        List of text chunks from each slide
    """
    try:
        prs = Presentation(io.BytesIO(file_content))
        text_chunks = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from all shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            # Clean and combine text
            combined_text = "\n".join(slide_text)
            cleaned_text = clean_text(combined_text)
            
            if cleaned_text:
                text_chunks.append(cleaned_text)
        
        return text_chunks
    except Exception as e:
        raise Exception(f"Error extracting text from PPTX: {str(e)}")


def clean_text(text: str) -> str:
    """
    Clean extracted text by removing extra whitespace and formatting
    
    Args:
        text: Raw extracted text
    
    Returns:
        Cleaned text
    """
    if not text:
        return ""
    
    # Remove excessive whitespace
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        line = line.strip()
        if line:
            cleaned_lines.append(line)
    
    return "\n".join(cleaned_lines)


def extract_text_by_type(file_content: bytes, file_type: str) -> List[str]:
    """
    Extract text based on file type
    
    Args:
        file_content: File content as bytes
        file_type: File MIME type (e.g., 'application/pdf', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
    
    Returns:
        List of text chunks
    """
    if file_type == 'application/pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return extract_text_from_pptx(file_content)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

